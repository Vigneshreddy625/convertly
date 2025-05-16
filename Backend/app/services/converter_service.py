import os
import subprocess
import logging
import pandas as pd
from pdf2docx import Converter as PDFToDocxConverter
from tabula import read_pdf
from pptx import Presentation
from PIL import Image
import fitz  # PyMuPDF
import tempfile

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def run_libreoffice_conversion(input_path: str, output_format: str, output_dir: str) -> bool:
    """Run LibreOffice to convert documents."""
    try:
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", output_format, 
            input_path, "--outdir", output_dir
        ], check=True, timeout=120)  # Added timeout
        return True
    except subprocess.CalledProcessError as e:
        logger.error(f"LibreOffice conversion failed: {e}")
        return False
    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return False

def word_to_pdf(word_path: str, output_path: str) -> bool:
    """Convert Word document to PDF."""
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)
    
    success = run_libreoffice_conversion(word_path, "pdf", output_dir)
    converted = os.path.join(output_dir, os.path.splitext(os.path.basename(word_path))[0] + ".pdf")
    
    if success and os.path.exists(converted):
        if converted != output_path:
            os.rename(converted, output_path)
        return True
    return False

def powerpoint_to_pdf(ppt_path: str, output_path: str) -> bool:
    """Convert PowerPoint to PDF."""
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)
    
    success = run_libreoffice_conversion(ppt_path, "pdf", output_dir)
    converted = os.path.join(output_dir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    
    if success and os.path.exists(converted):
        if converted != output_path:
            os.rename(converted, output_path)
        return True
    return False

def excel_to_pdf(excel_path: str, output_path: str) -> bool:
    """Convert Excel to PDF."""
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)
    
    success = run_libreoffice_conversion(excel_path, "pdf", output_dir)
    converted = os.path.join(output_dir, os.path.splitext(os.path.basename(excel_path))[0] + ".pdf")
    
    if success and os.path.exists(converted):
        if converted != output_path:
            os.rename(converted, output_path)
        return True
    return False

def pdf_to_word(pdf_path: str, output_path: str) -> bool:
    """Convert PDF to Word document."""
    try:
        output_dir = os.path.dirname(output_path)
        os.makedirs(output_dir, exist_ok=True)
        
        converter = PDFToDocxConverter(pdf_path)
        converter.convert(output_path)
        converter.close()
        return os.path.exists(output_path)
    except Exception as e:
        logger.error(f"PDF to Word conversion failed: {str(e)}")
        return False

def pdf_to_excel(pdf_path: str, output_path: str) -> bool:
    """Convert PDF to Excel, extracting tables."""
    try:
        output_dir = os.path.dirname(output_path)
        os.makedirs(output_dir, exist_ok=True)
        
        # Extract tables from PDF
        tables = read_pdf(pdf_path, pages='all', multiple_tables=True, lattice=True)
        
        if not tables:
            logger.warning(f"No tables found in {pdf_path}")
            # Create an empty Excel file with a note
            with pd.ExcelWriter(output_path) as writer:
                pd.DataFrame({"Note": ["No tables found in the PDF"]}).to_excel(
                    writer, sheet_name='Info', index=False
                )
            return os.path.exists(output_path)
        
        # Write tables to Excel
        with pd.ExcelWriter(output_path) as writer:
            for i, table in enumerate(tables):
                sheet_name = f'Table{i+1}'
                # Clean up the table data if needed
                table.to_excel(writer, sheet_name=sheet_name, index=False)
                
        return os.path.exists(output_path)
    except Exception as e:
        logger.error(f"PDF to Excel conversion failed: {str(e)}")
        return False

def pdf_to_powerpoint(pdf_path: str, output_path: str) -> bool:
    """Convert PDF to PowerPoint - create a slide for each page."""
    try:
        output_dir = os.path.dirname(output_path)
        os.makedirs(output_dir, exist_ok=True)
        
        # Create presentation
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Blank layout
        
        # Open PDF
        pdf_document = fitz.open(pdf_path)
        
        # Create a temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            # For each page in the PDF
            for page_num in range(len(pdf_document)):
                # Get the page
                page = pdf_document[page_num]
                
                # Add a slide
                slide = prs.slides.add_slide(slide_layout)
                
                # Convert page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Increase resolution
                img_path = os.path.join(temp_dir, f"page_{page_num+1}.png")
                pix.save(img_path)
                
                # Add image to slide
                with open(img_path, "rb") as img_file:
                    # Calculate the slide dimensions
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # Add image to fill most of the slide
                    left = slide_width * 0.05
                    top = slide_height * 0.05
                    width = slide_width * 0.9
                    height = slide_height * 0.9
                    
                    pic = slide.shapes.add_picture(
                        img_file, left, top, width=width, height=height
                    )
        
        # Save presentation
        prs.save(output_path)
        return os.path.exists(output_path)
    except Exception as e:
        logger.error(f"PDF to PowerPoint conversion failed: {str(e)}")
        return False
